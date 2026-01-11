import streamlit as st
from groq import Groq
from pypdf import PdfReader
import fitz  # PyMuPDF
import tempfile
import os
import re
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# =========================
# UI / CONFIG
# =========================
st.set_page_config(page_title="PPT Médico Pro", page_icon="🏥", layout="centered")
st.title("🏥 Generador Médico (Pro) — Slides + Figuras + Algoritmo")

if "GROQ_API_KEY" not in st.secrets:
    st.error("❌ Falta la GROQ_API_KEY en Secrets")
    st.stop()

client = Groq(api_key=st.secrets["GROQ_API_KEY"])

with st.sidebar:
    st.header("Opciones")

    num_content_slides = st.slider(
        "Nº de diapositivas de contenido (usuario)",
        min_value=4,
        max_value=60,
        value=12,
        step=1,
    )

    include_agenda = st.checkbox("Incluir Agenda", True)
    include_flow = st.checkbox("Incluir 1 slide de algoritmo/flujo", True)

    st.subheader("Figuras (imágenes desde el PDF)")
    include_figures = st.checkbox("Incluir figuras", True)
    max_fig_slides = st.slider("Cantidad de figuras", 0, 8, 3, disabled=not include_figures)

    figure_pages_text = st.text_input(
        "Páginas para figuras (opcional)",
        value="",
        help="Ejemplos: 3,5,7-9  (si lo dejas vacío, auto-selecciona)",
    )

    logo_file = st.file_uploader("Logo (opcional)", type=["png", "jpg", "jpeg"])

uploaded_file = st.file_uploader("Sube tu PDF médico", type="pdf")


# =========================
# STYLE HELPERS
# =========================
PRIMARY = RGBColor(0, 51, 102)       # Azul clínico
ACCENT = RGBColor(0, 102, 204)       # Azul acento
DARK = RGBColor(35, 35, 35)
LIGHT_BG = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(120, 120, 120)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_top_bar(slide, color=PRIMARY, height=Inches(0.20)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_footer(slide, left_text="", right_text=""):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(7.16), Inches(11.93), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 225, 230)
    line.line.fill.background()

    box_l = slide.shapes.add_textbox(Inches(0.7), Inches(7.20), Inches(8.2), Inches(0.3))
    tf = box_l.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = left_text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED

    box_r = slide.shapes.add_textbox(Inches(9.3), Inches(7.20), Inches(3.3), Inches(0.3))
    tf = box_r.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = right_text
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_logo(slide, logo_path):
    if not logo_path:
        return
    try:
        slide.shapes.add_picture(logo_path, Inches(11.75), Inches(0.35), height=Inches(0.6))
    except Exception:
        pass


def add_title(slide, title):
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.55), Inches(11.2), Inches(1.0))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    return tb


def add_key_message(slide, text):
    if not text:
        return
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.55), Inches(3.9), Inches(1.35))
    shp.fill.solid()
    shp.fill.fore_color.rgb = LIGHT_BG
    shp.line.color.rgb = RGBColor(210, 215, 220)
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER


def add_bullets(slide, bullets):
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(4.9))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.line_spacing = 1.15


def safe_filename(s):
    return re.sub(r"[^a-zA-Z0-9._-]+", "_", s)


# =========================
# PARSING HELPERS
# =========================
def parse_pages_input(text, max_pages):
    """
    '3,5,7-9' -> [2,4,6,7,8]  (0-based)
    """
    text = (text or "").strip()
    if not text:
        return []
    parts = [p.strip() for p in text.split(",") if p.strip()]
    out = []
    for p in parts:
        if "-" in p:
            a, b = p.split("-", 1)
            a = int(a.strip())
            b = int(b.strip())
            for x in range(a, b + 1):
                if 1 <= x <= max_pages:
                    out.append(x - 1)
        else:
            x = int(p)
            if 1 <= x <= max_pages:
                out.append(x - 1)
    # unique preserve order
    seen = set()
    res = []
    for x in out:
        if x not in seen:
            seen.add(x)
            res.append(x)
    return res


def parse_slides(llm_text):
    blocks = [b.strip() for b in llm_text.split("---SLIDE---") if b.strip()]
    slides = []
    for b in blocks:
        title = ""
        key = ""
        bullets = []
        mode = None
        for line in b.splitlines():
            line = line.strip()
            if not line:
                continue
            if line.startswith("TÍTULO:"):
                title = line.replace("TÍTULO:", "").strip()
                mode = None
            elif line.startswith("MENSAJE_CLAVE:"):
                key = line.replace("MENSAJE_CLAVE:", "").strip()
                mode = None
            elif line.startswith("CONTENIDO:"):
                mode = "bullets"
            elif mode == "bullets" and (line.startswith("•") or line.startswith("-")):
                bullets.append(line.lstrip("•-").strip())
        if title and bullets:
            slides.append({"title": title, "key": key, "bullets": bullets})
    return slides


# =========================
# GROQ (chunking)
# =========================
def gen_slide_batch(pdf_text, batch_n, batch_idx, total_n):
    prompt = f"""
Eres editor médico y diseñador de presentaciones (estilo congreso / sesión clínica).
Genera EXACTAMENTE {batch_n} diapositivas de CONTENIDO.

FORMATO OBLIGATORIO para cada una:
---SLIDE---
TÍTULO: <máx 8 palabras>
MENSAJE_CLAVE: <1 frase, máx 18 palabras>
CONTENIDO:
• <bullet 1, máx 14 palabras>
• <bullet 2, máx 14 palabras>
• <bullet 3, máx 14 palabras>

REGLAS:
- No introducciones fuera del formato.
- No inventes cifras.
- Sé clínico: criterios, riesgos, outcomes, diagnóstico, tratamiento.
- Evita redundancias.

Contexto: estás generando el lote {batch_idx+1} para completar {total_n} diapositivas.

TEXTO:
{pdf_text[:22000]}
""".strip()

    completion = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": "Responde estrictamente en el formato pedido."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.35,
        max_tokens=4500,
    )
    return completion.choices[0].message.content


def gen_all_slides(pdf_text, total_n, batch_size=10):
    all_text = []
    batches = (total_n + batch_size - 1) // batch_size
    for bi in range(batches):
        n = batch_size if bi < batches - 1 else (total_n - batch_size * (batches - 1))
        all_text.append(gen_slide_batch(pdf_text, n, bi, total_n))
    return "\n".join(all_text)


def gen_flow_steps(pdf_text):
    prompt = f"""
Crea un algoritmo clínico en 5–6 pasos máximo basado en el texto.
Devuelve SOLO una lista numerada:

1. ...
2. ...
3. ...

TEXTO:
{pdf_text[:18000]}
""".strip()

    completion = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": "Sé clínico, concreto y accionable."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.25,
        max_tokens=700,
    )
    raw = completion.choices[0].message.content.strip()
    steps = []
    for line in raw.splitlines():
        m = re.match(r"^\s*\d+\.\s+(.*)$", line.strip())
        if m:
            steps.append(m.group(1).strip())
    return steps[:6]


# =========================
# PDF -> PNG (figures)
# =========================
def render_page_png(pdf_bytes, page_index, zoom=2.0):
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    page = doc[page_index]
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    pix.save(tmp.name)
    doc.close()
    return tmp.name


def auto_pick_pages_for_figures(reader, max_k):
    """
    Auto: elige páginas donde aparezcan palabras tipo Figura/Tabla.
    Si no encuentra, usa las primeras páginas.
    """
    hits = []
    for i, page in enumerate(reader.pages):
        t = (page.extract_text() or "").lower()
        score = 0
        for kw in ["figura", "figure", "tabla", "table"]:
            if kw in t:
                score += 1
        if score > 0:
            hits.append((score, i))
    hits.sort(reverse=True, key=lambda x: x[0])
    pages = [i for _, i in hits][:max_k]
    if not pages:
        pages = list(range(min(max_k, len(reader.pages))))
    return pages


# =========================
# BUILD PPT
# =========================
def add_cover(prs, doc_name, today, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT_BG)
    add_top_bar(s, PRIMARY, Inches(0.22))
    add_logo(s, logo_path)

    tbox = s.shapes.add_textbox(Inches(1.2), Inches(2.15), Inches(11.0), Inches(2.2))
    tf = tbox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Resumen Clínico"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = doc_name
    p.font.size = Pt(20)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = f"Generado automáticamente · {today}"
    p.font.size = Pt(14)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER


def add_agenda(prs, doc_name, titles, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_top_bar(s)
    add_logo(s, logo_path)
    add_footer(s, left_text=doc_name, right_text="Agenda")
    add_title(s, "Agenda")

    # bullets (compacto)
    box = s.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.2))
    tf = box.text_frame
    tf.clear()
    for i, t in enumerate(titles[:12]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {t}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_content_slide(prs, doc_name, idx, total, item, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_top_bar(s)
    add_logo(s, logo_path)
    add_footer(s, left_text=doc_name, right_text=f"{idx}/{total}")

    add_title(s, item["title"])
    add_key_message(s, item.get("key", ""))
    add_bullets(s, item["bullets"])


def add_flow_slide(prs, doc_name, steps, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_top_bar(s, ACCENT)
    add_logo(s, logo_path)
    add_footer(s, left_text=doc_name, right_text="Algoritmo")

    add_title(s, "Algoritmo / Flujo clínico (resumen)")

    # vertical flow
    left = Inches(1.4)
    top = Inches(1.75)
    w = Inches(10.6)
    h = Inches(0.75)
    gap = Inches(0.15)

    count = min(len(steps), 6)
    for i in range(count):
        y = top + i * (h + gap)
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = LIGHT_BG
        shp.line.color.rgb = RGBColor(210, 215, 220)
        tf = shp.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {steps[i]}"
        p.font.size = Pt(16)
        p.font.bold = True if i == 0 else False
        p.font.color.rgb = PRIMARY

        if i < count - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.35), y + h, Inches(0.55), gap)
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT
            arr.line.fill.background()


def add_figure_slide(prs, doc_name, page_no_1based, png_path, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)
    add_top_bar(s, PRIMARY)
    add_logo(s, logo_path)
    add_footer(s, left_text=doc_name, right_text=f"Figura (pág. {page_no_1based})")

    add_title(s, f"Figura / Tabla (pág. {page_no_1based})")

    # Inserta imagen rendereada (página completa)
    # Ajuste simple: ancho fijo, deja margen
    s.shapes.add_picture(png_path, Inches(0.85), Inches(1.65), width=Inches(11.63))


def add_closing(prs, logo_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PRIMARY)
    add_logo(s, logo_path)

    box = s.shapes.add_textbox(Inches(2.0), Inches(2.8), Inches(9.3), Inches(1.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Gracias"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "¿Preguntas?"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


# =========================
# MAIN
# =========================
if uploaded_file:
    st.success(f"✅ Archivo cargado: {uploaded_file.name}")

    if st.button("🚀 Generar PowerPoint Pro"):
        logo_path = None
        tmp_pngs = []

        try:
            pdf_bytes = uploaded_file.getvalue()
            doc_name = uploaded_file.name.replace(".pdf", "")
            today = date.today().isoformat()

            # logo opcional
            if logo_file is not None:
                t = tempfile.NamedTemporaryFile(delete=False, suffix="." + logo_file.name.split(".")[-1])
                t.write(logo_file.getvalue())
                t.close()
                logo_path = t.name

            with st.spinner("📄 Extrayendo texto del PDF..."):
                reader = PdfReader(uploaded_file)
                pdf_text = ""
                for page in reader.pages:
                    pdf_text += (page.extract_text() or "") + "\n"

                if len(pdf_text.strip()) < 200:
                    st.warning("⚠️ El PDF tiene poco texto extraíble (puede ser escaneado). El contenido puede salir pobre.")

            with st.spinner("🧠 Generando diapositivas (IA) en lotes..."):
                llm_text = gen_all_slides(pdf_text, num_content_slides, batch_size=10)
                items = parse_slides(llm_text)

                if len(items) == 0:
                    st.error("La IA no devolvió el formato esperado. Dime y endurecemos el prompt.")
                    st.stop()

                # Si vienen más/menos, ajusta al número pedido (por seguridad)
                items = items[:num_content_slides]

            flow_steps = []
            if include_flow:
                with st.spinner("🧭 Generando algoritmo/flujo clínico..."):
                    flow_steps = gen_flow_steps(pdf_text)

            # FIGURAS: render de páginas (SIEMPRE funciona aunque sea vector)
            figure_pages = []
            if include_figures and max_fig_slides > 0:
                total_pages = len(reader.pages)
                chosen = parse_pages_input(figure_pages_text, total_pages)
                if not chosen:
                    chosen = auto_pick_pages_for_figures(reader, max_fig_slides)
                figure_pages = chosen[:max_fig_slides]

                with st.spinner("🖼️ Renderizando páginas a imágenes..."):
                    for pi in figure_pages:
                        png = render_page_png(pdf_bytes, pi, zoom=2.0)
                        tmp_pngs.append((pi, png))

            with st.spinner("🎨 Armando PowerPoint..."):
                prs = Presentation()
                prs.slide_width = SLIDE_W
                prs.slide_height = SLIDE_H

                add_cover(prs, doc_name, today, logo_path)

                if include_agenda:
                    add_agenda(prs, doc_name, [it["title"] for it in items], logo_path)

                for i, it in enumerate(items, start=1):
                    add_content_slide(prs, doc_name, i, len(items), it, logo_path)

                if include_flow and len(flow_steps) >= 3:
                    add_flow_slide(prs, doc_name, flow_steps, logo_path)

                for (pi, png) in tmp_pngs:
                    add_figure_slide(prs, doc_name, pi + 1, png, logo_path)

                add_closing(prs, logo_path)

            out_name = f"presentacion_medica_pro_{safe_filename(doc_name)}.pptx"
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                prs.save(tmp.name)
                with open(tmp.name, "rb") as f:
                    st.success(f"✅ OK. Generadas {len(prs.slides)} diapositivas (incluye portada/agenda/figuras).")
                    st.download_button(
                        "📥 Descargar PowerPoint",
                        data=f,
                        file_name=out_name,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )

            # limpiar tmp ppt
            try:
                os.unlink(tmp.name)
            except Exception:
                pass

        except Exception as e:
            st.error(f"Error: {e}")

        finally:
            # limpiar logo
            if logo_path:
                try:
                    os.unlink(logo_path)
                except Exception:
                    pass
            # limpiar pngs
            for _, p in tmp_pngs:
                try:
                    os.unlink(p)
                except Exception:
                    pass
else:
    st.info("⬆️ Sube un PDF para comenzar.")
